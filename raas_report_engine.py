"""raas_report_engine.py — 정적 보고서(PPT) 생성 엔진 (Phase 3 ⑥).

matplotlib로 차트를 PNG로 렌더 + python-pptx로 한 장짜리 보고서 작성.
스토리라인 산출물(CP 국장 1장 PPT 등)이 이 모듈로 변환된다.

진입점:
    build_storyline_export(role, output_format_id, slots_visited, context)
        → {ok, filename, content_type, bytes}

설계:
    - 9종 chart_type별 matplotlib 렌더러 (현재 line / bar 우선, 나머지 폴백)
    - 한국어 폰트 자동 탐지 (Malgun Gothic 우선)
    - PPT 1장 = 제목 + 본문 텍스트 + 차트 1~2개 + 푸터(생성일·역할)
    - 채팅 답변의 변수 치환 결과를 받아 그대로 옮김
"""
from __future__ import annotations

import io
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

# matplotlib는 GUI 백엔드 불필요
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
EXPORTS_DIR = ROOT / "data" / "exports"
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)


# ─── 한글 폰트 자동 탐지 ────────────────────────────────────────────
def _detect_korean_font() -> str:
    """OS별 한글 폰트 우선순위 — Windows: Malgun Gothic / 그 외: 사용 가능한 첫 한글 폰트."""
    candidates = [
        "Malgun Gothic", "맑은 고딕",
        "NanumGothic", "나눔고딕",
        "AppleGothic", "Apple SD Gothic Neo",
        "Noto Sans CJK KR", "Noto Sans KR",
    ]
    available = {f.name for f in font_manager.fontManager.ttflist}
    for c in candidates:
        if c in available:
            return c
    return "DejaVu Sans"


KOREAN_FONT = _detect_korean_font()
matplotlib.rcParams["font.family"] = KOREAN_FONT
matplotlib.rcParams["axes.unicode_minus"] = False


# ─── 색상 팔레트 (ECharts 프론트와 일치) ──────────────────────────
PALETTE = [
    "#e74c3c", "#3498db", "#16a085", "#f39c12",
    "#9b59b6", "#1abc9c", "#e67e22", "#34495e",
]


# ─── matplotlib 차트 렌더러 (PNG 반환) ───────────────────────────
def render_chart_png(chart: dict, width: int = 8, height: int = 4, dpi: int = 150) -> bytes:
    """chart_data dict → PNG 바이트.

    지원: timeseries / timeseries_multi / timeseries_dual / comparison
    그 외: 빈 안내 이미지 반환.
    """
    fig, ax = plt.subplots(figsize=(width, height), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    ctype = chart.get("type") or ""
    chart_kind = chart.get("chart_type") or ""

    try:
        if ctype == "timeseries" and chart.get("points"):
            _render_single_line(ax, chart)
        elif ctype == "timeseries_multi" and chart.get("series"):
            _render_multi_line(ax, chart, dual=False)
        elif ctype == "timeseries_dual" and chart.get("series"):
            _render_multi_line(ax, chart, dual=True)
        elif ctype == "comparison" and chart.get("items"):
            _render_bar(ax, chart, kind=chart_kind or "rank")
        elif chart_kind in {"bar_rank", "bar_delta"} and chart.get("items"):
            _render_bar(ax, chart, kind=chart_kind.replace("bar_", ""))
        elif chart.get("series"):
            _render_multi_line(ax, chart, dual=False)
        elif chart.get("points"):
            _render_single_line(ax, chart)
        elif chart.get("items"):
            _render_bar(ax, chart, kind="rank")
        else:
            ax.text(0.5, 0.5, "지원하지 않는 차트 형태", ha="center", va="center",
                    color="#888888", fontsize=12, transform=ax.transAxes)
            ax.axis("off")
    except Exception as e:
        ax.clear()
        ax.text(0.5, 0.5, f"차트 렌더 실패: {e}", ha="center", va="center",
                color="#c0392b", fontsize=10, transform=ax.transAxes)
        ax.axis("off")

    title = chart.get("title")
    if title:
        ax.set_title(title, fontsize=13, fontweight="bold", color="#222", pad=10)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _setup_axes(ax) -> None:
    """공통 축 디자인."""
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    ax.grid(True, axis="y", linestyle="--", alpha=0.4, color="#cccccc")
    ax.tick_params(colors="#666666", labelsize=9)


def _render_single_line(ax, chart: dict) -> None:
    pts = chart.get("points") or []
    dates = [p.get("date", "") for p in pts]
    values = [p.get("value", 0) for p in pts]
    color = chart.get("color") or PALETTE[1]
    label = chart.get("label") or chart.get("metric") or "값"
    ax.plot(dates, values, color=color, marker="o", markersize=4, linewidth=2, label=label)
    _setup_axes(ax)
    _rotate_xlabels(ax, dates)
    if chart.get("unit") == "%":
        ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.legend(loc="upper left", frameon=False, fontsize=9)


def _render_multi_line(ax, chart: dict, dual: bool = False) -> None:
    series = chart.get("series") or []
    if not series:
        return
    dates = [p.get("date", "") for p in (series[0].get("points") or [])]
    units = [s.get("unit", "") for s in series]
    same_unit = len(set(units)) == 1

    if dual and not same_unit and len(series) >= 2:
        # 이중 축
        ax2 = ax.twinx()
        for i, s in enumerate(series[:2]):
            target = ax if i == 0 else ax2
            color = s.get("color") or PALETTE[i % len(PALETTE)]
            vals = [p.get("value", 0) for p in (s.get("points") or [])]
            target.plot(dates, vals, color=color, marker="o", markersize=4,
                        linewidth=2, label=s.get("label", ""))
            target.tick_params(axis="y", colors=color)
            target.spines["top"].set_visible(False)
            if i == 0:
                target.spines["right"].set_visible(False)
            else:
                target.spines["left"].set_visible(False)
        ax.spines["bottom"].set_color("#cccccc")
        ax.grid(True, axis="y", linestyle="--", alpha=0.4, color="#cccccc")
        ax.tick_params(axis="x", colors="#666666", labelsize=9)
        _rotate_xlabels(ax, dates)
        lines1, labels1 = ax.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax.legend(lines1 + lines2, labels1 + labels2, loc="upper left",
                  frameon=False, fontsize=9)
    else:
        for i, s in enumerate(series):
            color = s.get("color") or PALETTE[i % len(PALETTE)]
            vals = [p.get("value", 0) for p in (s.get("points") or [])]
            ax.plot(dates, vals, color=color, marker="o", markersize=4,
                    linewidth=2, label=s.get("label", f"시리즈 {i+1}"))
        _setup_axes(ax)
        _rotate_xlabels(ax, dates)
        if same_unit and units[0] == "%":
            ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
        ax.legend(loc="upper left", frameon=False, fontsize=9)


def _render_bar(ax, chart: dict, kind: str = "rank") -> None:
    items = chart.get("items") or []
    if kind == "rank":
        items = sorted(items, key=lambda it: -(it.get("value") or 0))
    labels = [it.get("label", "") for it in items]
    values = [it.get("value", 0) for it in items]
    if kind == "delta":
        colors = ["#16a085" if v >= 0 else "#e74c3c" for v in values]
    else:
        colors = [it.get("color") or PALETTE[i % len(PALETTE)] for i, it in enumerate(items)]
    bars = ax.bar(labels, values, color=colors, width=0.6)
    _setup_axes(ax)
    # 값 라벨
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height(),
                f"{v:,}" if chart.get("unit") != "%" else f"{v:.1f}%",
                ha="center", va="bottom", fontsize=9, color="#222")
    if chart.get("unit") == "%":
        ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    if any(len(str(l)) > 4 for l in labels):
        plt.setp(ax.get_xticklabels(), rotation=30, ha="right")


def _rotate_xlabels(ax, labels: list) -> None:
    if len(labels) > 20:
        plt.setp(ax.get_xticklabels(), rotation=30, ha="right")
    # 너무 많은 라벨이면 sub-sampling
    if len(labels) > 40:
        step = max(1, len(labels) // 12)
        for i, t in enumerate(ax.get_xticklabels()):
            if i % step != 0:
                t.set_visible(False)


# ─── PPT 생성기 ────────────────────────────────────────────────────
def _build_one_pager_ppt(
    title: str,
    body_text: str,
    chart_pngs: list[bytes],
    footer: str,
) -> bytes:
    """한 장짜리 보고서 PPT 작성.

    레이아웃:
        [제목]                          (상단)
        [본문 텍스트]      [차트 1]
                          [차트 2]
        [푸터]                          (하단)
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    run.font.name = KOREAN_FONT

    # 본문 텍스트
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.0))
    btf = body_box.text_frame
    btf.word_wrap = True
    for ln in body_text.splitlines():
        para = btf.add_paragraph() if btf.paragraphs[0].text else btf.paragraphs[0]
        para.text = ln
        for r in para.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            r.font.name = KOREAN_FONT

    # 차트 이미지 (최대 2개)
    chart_left = Inches(5.0)
    chart_width = Inches(4.8)
    chart_height = Inches(2.5)
    for i, png in enumerate(chart_pngs[:2]):
        top = Inches(1.2 + i * 2.7)
        slide.shapes.add_picture(io.BytesIO(png), chart_left, top,
                                 width=chart_width, height=chart_height)

    # 푸터
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9.0), Inches(0.4))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = footer
    fp.alignment = PP_ALIGN.RIGHT
    frun = fp.runs[0]
    frun.font.size = Pt(9)
    frun.font.color.rgb = RGBColor(0x8B, 0x95, 0xA5)
    frun.font.italic = True
    frun.font.name = KOREAN_FONT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── 스토리라인 산출물 통합 진입점 ────────────────────────────────
def build_storyline_export(
    role: str,
    output_format_id: str,
    slots_visited: list,
    context: Optional[dict] = None,
) -> dict:
    """스토리라인 5_closing의 '내보내기' 칩 → 실제 파일 생성.

    Args:
        role: 'CP' / 'PD_PROGRAM' / 'PD_SCHEDULE' (대소문자 무관)
        output_format_id: cp.json outputsTo.id (예: 'director_one_pager_ppt')
        slots_visited: [{slot, answer, fallback_used, ...}, ...]
        context: 마지막 advance의 context_out (변수 치환용)

    Returns:
        {ok, filename, content_type, bytes}
    """
    context = context or {}
    role_upper = (role or "").upper()
    fmt_id = output_format_id or ""

    # 본문 텍스트 — 방문한 슬롯들의 answer 모음
    sections = []
    for s in slots_visited:
        if not s: continue
        slot_name = s.get("slot_name") or s.get("slot") or ""
        answer = s.get("answer") or ""
        if not answer.strip(): continue
        # 마크다운 **bold** → 일반 텍스트
        clean_answer = re.sub(r"\*\*(.+?)\*\*", r"\1", answer)
        sections.append(f"■ {slot_name}\n{clean_answer}\n")
    body_text = "\n".join(sections) or "(분석 결과가 없습니다.)"

    # 차트 PNG — slots_visited의 chart_data가 있으면 렌더
    chart_pngs: list[bytes] = []
    for s in slots_visited:
        chart_data = s.get("chart_data") if isinstance(s, dict) else None
        if isinstance(chart_data, dict):
            try:
                chart_pngs.append(render_chart_png(chart_data))
            except Exception as e:
                print(f"[report] chart render failed: {e}")

    # 제목 / 푸터 (산출물별 분기)
    today = datetime.now().strftime("%Y-%m-%d")
    user_name = context.get("user_name", "")
    title_map = {
        "director_one_pager_ppt": f"국장 보고용 — {context.get('channel_name','')} 주간 핵심 ({today})",
        "weekly_meeting_free_format": f"주간 회의 자료 ({today})",
        "business_director_mau_brief": f"사업국장 MAU 별도 보고 ({today})",
        "pd_dispatch_message": "PD 디스패치 메시지",
        "cp_one_pager": f"CP 보고용 1장 ({today})",
        "writer_host_share_kakao": "작가/진행자 공유 메시지",
        "scheduling_meeting_formal_report": f"편성 회의 보고서 ({today})",
        "director_summary_table": f"국장 보고용 요약표 ({today})",
        "listening_survey_side_by_side": f"청취율 조사 요약 ({today})",
        "comparison_graph_export": f"비교 그래프 ({today})",
        "graph_excel_download": f"그래프 데이터 ({today})",
        "png_save_mobile": f"모바일 공유 ({today})",
    }
    title = title_map.get(fmt_id, f"{fmt_id} ({today})")
    footer = f"{user_name} · {today} · RAAS 자동 생성"

    # 텍스트 출력물 (카톡 등)
    if fmt_id in {"pd_dispatch_message", "writer_host_share_kakao"}:
        text = f"{title}\n\n{body_text}\n\n— {footer}"
        return {
            "ok": True,
            "filename": f"{fmt_id}_{today}.txt",
            "content_type": "text/plain; charset=utf-8",
            "bytes": text.encode("utf-8"),
            "kind": "text",
        }

    # PPT 출력물 (대부분)
    ppt_bytes = _build_one_pager_ppt(title, body_text, chart_pngs, footer)
    return {
        "ok": True,
        "filename": f"{fmt_id}_{today}.pptx",
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "bytes": ppt_bytes,
        "kind": "ppt",
    }


# ─── 일회용 토큰 기반 다운로드 캐시 ────────────────────────────────
# /api/storyline/export가 응답에 다운로드 URL을 담아주고, 그 URL이 호출되면 파일 반환.
_EXPORT_CACHE: dict[str, dict] = {}


def store_for_download(result: dict) -> str:
    """build_storyline_export 결과를 캐시에 저장하고 토큰 반환."""
    import secrets
    token = secrets.token_urlsafe(16)
    _EXPORT_CACHE[token] = result
    # 캐시는 최대 32개만 보관 (FIFO)
    if len(_EXPORT_CACHE) > 32:
        oldest = next(iter(_EXPORT_CACHE))
        _EXPORT_CACHE.pop(oldest, None)
    return token


def fetch_for_download(token: str) -> Optional[dict]:
    """토큰으로 캐시된 export 결과를 1회 가져오기 (pop)."""
    return _EXPORT_CACHE.pop(token, None)


if __name__ == "__main__":
    # CLI 자가 테스트
    sample_chart = {
        "type": "timeseries_multi",
        "title": "테스트 채널별 DAU",
        "series": [
            {"label": "파워FM", "points": [{"date":"06-10","value":120000},{"date":"06-11","value":118000},
                                       {"date":"06-12","value":121000},{"date":"06-13","value":125000}]},
            {"label": "러브FM", "points": [{"date":"06-10","value":90000},{"date":"06-11","value":89000},
                                       {"date":"06-12","value":92000},{"date":"06-13","value":94000}]},
        ],
    }
    png = render_chart_png(sample_chart)
    out = ROOT / "data" / "exports" / "test_chart.png"
    out.write_bytes(png)
    print(f"chart png: {out} ({len(png)} bytes)")

    result = build_storyline_export(
        role="CP",
        output_format_id="director_one_pager_ppt",
        slots_visited=[
            {"slot": "1_anchor", "slot_name": "묶음 + 변화 크기",
             "answer": "파워FM의 **13개 프로그램** 중 **두시탈출 컬투쇼**가 -7.3% 변화.",
             "chart_data": sample_chart},
            {"slot": "2_cause", "slot_name": "원인 설명",
             "answer": "변화 원인 후보:\n1. 게스트 효과 감소\n2. 인접 흡수"},
            {"slot": "5_closing", "slot_name": "마무리"},
        ],
        context={"user_name": "박CP", "channel_name": "파워FM"},
    )
    out_ppt = ROOT / "data" / "exports" / result["filename"]
    out_ppt.write_bytes(result["bytes"])
    print(f"PPT: {out_ppt} ({len(result['bytes'])} bytes)")
